"""
Generate Saloon Pro CRM PowerPoint Presentation
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──
DARK = RGBColor(0x0F, 0x17, 0x2A)
DARK2 = RGBColor(0x1E, 0x29, 0x3B)
PRIMARY = RGBColor(0x0D, 0x94, 0x88)
PRIMARY_DARK = RGBColor(0x0F, 0x76, 0x6E)
PRIMARY_LIGHT = RGBColor(0x99, 0xF6, 0xE4)
ACCENT = RGBColor(0xF5, 0x9E, 0x0B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
LIGHT_GRAY = RGBColor(0xCB, 0xD5, 0xE1)
GREEN_CHECK = RGBColor(0x34, 0xD3, 0x99)
CARD_BG = RGBColor(0x16, 0x2A, 0x3A)

IMAGES_DIR = os.path.join(os.path.dirname(__file__), "images")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius is not None:
        shape.adjustments[0] = radius
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri", line_space=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_space != 1.0:
        p.line_spacing = Pt(font_size * line_space)
    return txBox


def add_multiline_box(slide, left, top, width, height, lines, default_size=16,
                      default_color=WHITE, font_name="Calibri", alignment=PP_ALIGN.LEFT):
    """lines = list of (text, font_size, color, bold)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        text = line_data[0]
        size = line_data[1] if len(line_data) > 1 else default_size
        color = line_data[2] if len(line_data) > 2 else default_color
        bold = line_data[3] if len(line_data) > 3 else False
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(4)
    return txBox


def add_check_list(slide, left, top, width, items, font_size=15, spacing=6):
    height = Inches(len(items) * 0.38)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run1 = p.add_run()
        run1.text = "\u2713  "
        run1.font.size = Pt(font_size)
        run1.font.color.rgb = GREEN_CHECK
        run1.font.bold = True
        run1.font.name = "Calibri"
        run2 = p.add_run()
        run2.text = item
        run2.font.size = Pt(font_size)
        run2.font.color.rgb = LIGHT_GRAY
        run2.font.bold = False
        run2.font.name = "Calibri"
        p.space_after = Pt(spacing)
        p.space_before = Pt(0)
    return txBox


def add_image_safe(slide, path, left, top, width=None, height=None):
    full = os.path.join(IMAGES_DIR, path)
    if os.path.exists(full):
        if width and height:
            slide.shapes.add_picture(full, left, top, width, height)
        elif width:
            slide.shapes.add_picture(full, left, top, width=width)
        elif height:
            slide.shapes.add_picture(full, left, top, height=height)
        else:
            slide.shapes.add_picture(full, left, top)


def add_stat_card(slide, left, top, number, label):
    w, h = Inches(2.6), Inches(1.7)
    shape = add_shape_bg(slide, left, top, w, h, CARD_BG, radius=0.08)
    add_text_box(slide, left, top + Inches(0.25), w, Inches(0.7),
                 number, font_size=40, color=PRIMARY_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + Inches(1.05), w, Inches(0.4),
                 label, font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)


def add_module_card(slide, left, top, icon, title, desc):
    w, h = Inches(3.8), Inches(2.0)
    shape = add_shape_bg(slide, left, top, w, h, CARD_BG, radius=0.06)
    add_text_box(slide, left + Inches(0.25), top + Inches(0.2), Inches(0.5), Inches(0.5),
                 icon, font_size=24, color=PRIMARY_LIGHT)
    add_text_box(slide, left + Inches(0.25), top + Inches(0.65), w - Inches(0.5), Inches(0.4),
                 title, font_size=16, color=WHITE, bold=True)
    add_text_box(slide, left + Inches(0.25), top + Inches(1.05), w - Inches(0.5), Inches(0.85),
                 desc, font_size=11, color=GRAY, line_space=1.3)


def add_badge(slide, left, top, text, bg_color=CARD_BG, text_color=PRIMARY_LIGHT):
    w = Inches(len(text) * 0.09 + 0.3)
    shape = add_shape_bg(slide, left, top, w, Inches(0.32), bg_color, radius=0.15)
    add_text_box(slide, left, top + Inches(0.02), w, Inches(0.28),
                 text, font_size=10, color=text_color, bold=True, alignment=PP_ALIGN.CENTER)
    return w


# ══════════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK)

# Decorative accent bar
add_shape_bg(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, PRIMARY)

# Badge
add_shape_bg(slide, Inches(0.8), Inches(1.0), Inches(3.0), Inches(0.38), CARD_BG, radius=0.3)
add_text_box(slide, Inches(0.8), Inches(1.03), Inches(3.0), Inches(0.35),
             "\u25CF  Cloud-Powered SaaS Platform", font_size=11, color=PRIMARY_LIGHT, alignment=PP_ALIGN.CENTER)

# Title
add_text_box(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(1.8),
             "The Complete\nSalon Management\nSystem", font_size=44, color=WHITE, bold=True,
             font_name="Calibri", line_space=1.15)

# Subtitle
add_text_box(slide, Inches(0.8), Inches(3.8), Inches(5.2), Inches(1.0),
             "Enterprise-grade, multi-branch salon & spa management platform.\nBilling, CRM, appointments, inventory, staff, analytics \u2014 all in one.",
             font_size=16, color=GRAY, line_space=1.5)

# Stats row
stats = [("50+", "Modules"), ("450+", "API Endpoints"), ("25+", "Reports"), ("30+", "Data Models")]
for i, (num, lbl) in enumerate(stats):
    x = Inches(0.8 + i * 1.5)
    add_text_box(slide, x, Inches(5.2), Inches(1.3), Inches(0.5), num, font_size=28, color=WHITE, bold=True)
    add_text_box(slide, x, Inches(5.65), Inches(1.3), Inches(0.3), lbl, font_size=10, color=GRAY)

# Dashboard image
add_image_safe(slide, "dashboard-mockup.png", Inches(6.8), Inches(0.8), width=Inches(6.0))

print("  Slide 1: Title")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 2: PLATFORM OVERVIEW
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_text_box(slide, Inches(0), Inches(0.5), SLIDE_W, Inches(0.7),
             "Platform at a Glance", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(1.15), SLIDE_W, Inches(0.4),
             "A comprehensive solution built for modern salon businesses", font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

# Stat cards
card_x_start = Inches(0.85)
for i, (num, lbl) in enumerate([("50+", "Frontend Modules"), ("450+", "API Endpoints"),
                                 ("25+", "Business Reports"), ("30+", "Data Collections")]):
    add_stat_card(slide, card_x_start + Inches(i * 3.1), Inches(2.0), num, lbl)

# Capability pills
caps = [("\U0001F4B0", "Point of Sale"), ("\U0001F464", "CRM"), ("\U0001F4C5", "Appointments"),
        ("\U0001F465", "Staff & HR"), ("\U0001F4CA", "Analytics")]
for i, (icon, label) in enumerate(caps):
    x = Inches(0.85 + i * 2.5)
    shape = add_shape_bg(slide, x, Inches(4.4), Inches(2.2), Inches(1.1), CARD_BG, radius=0.08)
    add_text_box(slide, x, Inches(4.55), Inches(2.2), Inches(0.45), icon, font_size=28, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(5.0), Inches(2.2), Inches(0.35), label, font_size=13, color=WHITE,
                 bold=True, alignment=PP_ALIGN.CENTER)

print("  Slide 2: Overview")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 3: POS & BILLING
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_shape_bg(slide, Inches(0.6), Inches(0.6), Inches(1.6), Inches(0.35), CARD_BG, radius=0.3)
add_text_box(slide, Inches(0.6), Inches(0.62), Inches(1.6), Inches(0.32),
             "CORE MODULE", font_size=10, color=PRIMARY_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.6), Inches(1.2), Inches(5.5), Inches(0.8),
             "Point of Sale &\nSmart Billing", font_size=36, color=WHITE, bold=True, line_space=1.1)

add_text_box(slide, Inches(0.6), Inches(2.3), Inches(5.2), Inches(0.7),
             "Lightning-fast checkout built for high-volume salon operations.\nMulti-item bills, multiple payment methods, and professional GST invoices.",
             font_size=14, color=GRAY, line_space=1.5)

add_check_list(slide, Inches(0.6), Inches(3.3), Inches(5.5), [
    "Services, Packages, Products, Memberships & Prepaid",
    "Cash, Card, UPI, Wallet payment methods",
    "GST-compliant PDF invoice generation",
    "Real-time inventory deduction on sale",
    "Discount approval workflow with manager codes",
    "Share invoices via Email, WhatsApp, public link",
    "Bill history with deleted bills archive & audit trail",
])

add_image_safe(slide, "billing-mockup.png", Inches(6.8), Inches(0.7), width=Inches(6.0))

print("  Slide 3: POS")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 4: CRM
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_image_safe(slide, "customer-mockup.png", Inches(0.5), Inches(0.7), width=Inches(6.0))

add_shape_bg(slide, Inches(7.0), Inches(0.6), Inches(2.5), Inches(0.35), CARD_BG, radius=0.3)
add_text_box(slide, Inches(7.0), Inches(0.62), Inches(2.5), Inches(0.32),
             "CUSTOMER INTELLIGENCE", font_size=10, color=PRIMARY_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(7.0), Inches(1.2), Inches(5.8), Inches(0.8),
             "360\u00b0 Customer\nRelationship Management", font_size=36, color=WHITE, bold=True, line_space=1.1)

add_text_box(slide, Inches(7.0), Inches(2.3), Inches(5.8), Inches(0.7),
             "Build lasting relationships with complete customer profiles,\nlead tracking, feedback management, and lifecycle analytics.",
             font_size=14, color=GRAY, line_space=1.5)

add_check_list(slide, Inches(7.0), Inches(3.3), Inches(5.8), [
    "Full customer profile with visit & spend history",
    "Lead management with conversion tracking",
    "Missed enquiry follow-up system",
    "Public feedback page (no login required)",
    "Customer merge & deduplication with audit trail",
    "VIP / Regular / At-Risk lifecycle segmentation",
    "Referral program with attribution tracking",
])

print("  Slide 4: CRM")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 5: APPOINTMENTS
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_shape_bg(slide, Inches(0.6), Inches(0.6), Inches(1.5), Inches(0.35), CARD_BG, radius=0.3)
add_text_box(slide, Inches(0.6), Inches(0.62), Inches(1.5), Inches(0.32),
             "SCHEDULING", font_size=10, color=PRIMARY_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.6), Inches(1.2), Inches(5.5), Inches(0.8),
             "Visual Appointment\nCalendar", font_size=36, color=WHITE, bold=True, line_space=1.1)

add_text_box(slide, Inches(0.6), Inches(2.3), Inches(5.2), Inches(0.7),
             "A beautiful calendar system for managing bookings, staff\navailability, and appointment workflows across all branches.",
             font_size=14, color=GRAY, line_space=1.5)

add_check_list(slide, Inches(0.6), Inches(3.3), Inches(5.5), [
    "Weekly & monthly calendar views",
    "Staff assignment with real-time availability",
    "Confirmed / Completed / Cancelled / No-show status",
    "Dynamic available slot calculation",
    "Direct appointment-to-bill linking",
    "Multi-branch scheduling support",
    "Booking trends & cancellation analytics",
])

add_image_safe(slide, "appointment-mockup.png", Inches(6.8), Inches(0.7), width=Inches(6.0))

print("  Slide 5: Appointments")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 6: ALL MODULES
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_text_box(slide, Inches(0), Inches(0.3), SLIDE_W, Inches(0.6),
             "Complete Business Platform", font_size=34, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(0.9), SLIDE_W, Inches(0.35),
             "Every module works together for seamless salon operations", font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

modules = [
    ("\U0001F4B0", "Financial Management", "Cash register, expense tracking,\ntax slabs, GST compliance"),
    ("\U0001F4BC", "Membership & Prepaid", "Membership plans, prepaid balances,\nreal-time usage & expiry alerts"),
    ("\U0001F4E6", "Inventory Management", "Real-time stock, auto-deduction,\nlow stock alerts, SKU tracking"),
    ("\U0001F465", "Staff & HR", "Attendance, leave management,\nperformance, commission tracking"),
    ("\U0001F31F", "Service Catalog", "Service groups, combo packages,\npricing, branch-specific availability"),
    ("\U0001F512", "Discount Approvals", "Staff request, manager approve,\ngenerated codes & audit trail"),
    ("\U0001F4E3", "Marketing Campaigns", "Birthday campaigns, offers,\ncustomer targeting, WhatsApp"),
    ("\U0001F3E2", "Multi-Branch Ops", "Data isolation, temp assignments,\ncross-branch owner visibility"),
    ("\U0001F6E0\uFE0F", "Asset Management", "Fixed assets by branch, purchase\ndetails, status & reporting"),
]

for i, (icon, title, desc) in enumerate(modules):
    row = i // 3
    col = i % 3
    x = Inches(0.6 + col * 4.15)
    y = Inches(1.55 + row * 2.05)
    add_module_card(slide, x, y, icon, title, desc)

print("  Slide 6: Modules")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 7: REPORTS & ANALYTICS
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_image_safe(slide, "reports-mockup.png", Inches(0.5), Inches(0.7), width=Inches(6.2))

add_shape_bg(slide, Inches(7.2), Inches(0.6), Inches(2.2), Inches(0.35),
             RGBColor(0x30, 0x28, 0x10), radius=0.3)
add_text_box(slide, Inches(7.2), Inches(0.62), Inches(2.2), Inches(0.32),
             "BUSINESS INTELLIGENCE", font_size=10, color=RGBColor(0xFC, 0xD3, 0x4D),
             bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(7.2), Inches(1.2), Inches(5.5), Inches(0.6),
             "25+ Powerful Reports", font_size=36, color=WHITE, bold=True)

add_text_box(slide, Inches(7.2), Inches(2.0), Inches(5.5), Inches(0.7),
             "Data-driven decisions made easy. Real-time KPIs,\ntrend analysis, and deep-dive business insights.",
             font_size=14, color=GRAY, line_space=1.5)

# Report tags
reports = [
    "Revenue Dashboard", "Staff Performance", "Sales Insights", "Customer Lifecycle",
    "Growth Trends", "Client Value", "Expense Reports", "Inventory",
    "Incentives", "Membership", "Top Customers", "Payments",
    "Period Summary", "Service Analysis", "Top Items"
]

y_pos = Inches(3.0)
x_pos = Inches(7.2)
row_width = 0
for tag in reports:
    tag_w = Inches(len(tag) * 0.1 + 0.35)
    if row_width + tag_w.inches > 5.5:
        y_pos += Inches(0.42)
        x_pos = Inches(7.2)
        row_width = 0
    shape = add_shape_bg(slide, x_pos, y_pos, tag_w, Inches(0.33), CARD_BG, radius=0.2)
    add_text_box(slide, x_pos, y_pos + Inches(0.03), tag_w, Inches(0.28),
                 tag, font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    x_pos += tag_w + Inches(0.1)
    row_width += tag_w.inches + 0.1

print("  Slide 7: Reports")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 8: TECH STACK
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_text_box(slide, Inches(0), Inches(0.4), SLIDE_W, Inches(0.6),
             "Modern, Scalable Architecture", font_size=34, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(1.0), SLIDE_W, Inches(0.35),
             "Built with industry-leading technologies for performance and reliability",
             font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

# Three tech columns
tech_data = [
    ("\u2699\uFE0F", "Frontend", "React 18 + Vite 5", [
        "React 18", "Vite 5", "Ant Design 6", "Recharts",
        "Zustand", "React Query", "Framer Motion"
    ]),
    ("\U0001F5A5\uFE0F", "Backend", "Python / Flask 3", [
        "Flask 3", "MongoEngine", "PyJWT", "bcrypt",
        "ReportLab PDF", "Gunicorn", "flask-compress"
    ]),
    ("\u2601\uFE0F", "Infrastructure", "Cloud Deployment", [
        "MongoDB Atlas", "Google Cloud Run", "Docker",
        "Redis Cache", "gzip Compression"
    ]),
]

for i, (icon, title, subtitle, tags) in enumerate(tech_data):
    x = Inches(0.6 + i * 4.15)
    y = Inches(1.7)
    w = Inches(3.85)
    h = Inches(5.0)
    shape = add_shape_bg(slide, x, y, w, h, CARD_BG, radius=0.06)

    add_text_box(slide, x, y + Inches(0.3), w, Inches(0.5),
                 icon, font_size=32, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(0.85), w, Inches(0.4),
                 title, font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(1.25), w, Inches(0.3),
                 subtitle, font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Tags
    tag_y = y + Inches(1.75)
    for tag in tags:
        tw = Inches(2.4)
        tag_shape = add_shape_bg(slide, x + Inches(0.72), tag_y, tw, Inches(0.35),
                                  RGBColor(0x10, 0x3A, 0x36), radius=0.15)
        add_text_box(slide, x + Inches(0.72), tag_y + Inches(0.03), tw, Inches(0.3),
                     tag, font_size=12, color=PRIMARY_LIGHT, alignment=PP_ALIGN.CENTER, bold=True)
        tag_y += Inches(0.42)

print("  Slide 8: Tech Stack")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 9: ARCHITECTURE — WHAT POWERS WHAT
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_text_box(slide, Inches(0), Inches(0.3), SLIDE_W, Inches(0.55),
             "What Powers What", font_size=34, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(0.85), SLIDE_W, Inches(0.3),
             "Detailed breakdown of technologies and where they are used",
             font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

# Three columns
arch_cols = [
    ("\u2699\uFE0F", "Frontend", "React 18 + Vite 5", RGBColor(0x3B, 0x82, 0xF6), [
        ("Ant Design", "All UI: tables, forms, modals, pickers"),
        ("Recharts", "Dashboard charts, revenue, analytics"),
        ("React Query", "All API calls, caching, auto-refetch"),
        ("Zustand", "Sidebar, theme, global UI state"),
        ("Framer Motion", "Page transitions, modal animations"),
        ("Hook Form", "Customer, staff, service forms"),
        ("Day.js", "Calendar, filters, membership expiry"),
    ]),
    ("\U0001F5A5\uFE0F", "Backend", "Flask 3 + Gunicorn", PRIMARY, [
        ("Blueprints", "34 modules: /api/bills, /customers..."),
        ("MongoEngine", "30+ models: Bill, Customer, Staff..."),
        ("PyJWT", "Token auth on every protected route"),
        ("bcrypt", "Password hash for Staff/Manager/Owner"),
        ("ReportLab", "GST-compliant PDF invoice generation"),
        ("Aggregation", "Leaderboard, revenue, funnel, reports"),
        ("Redis/Cache", "Dashboard stats 5min TTL, report data"),
    ]),
    ("\U0001F5C3\uFE0F", "Database", "MongoDB Atlas \u2014 30+ Collections", ACCENT, [
        ("Core", "branches, customers, staffs, owners"),
        ("Catalog", "services, products, packages, groups"),
        ("Billing", "bills (embedded items), invoices, cash"),
        ("Membership", "plans, memberships, prepaid balances"),
        ("HR", "attendance, leaves, temp assignments"),
        ("CRM", "leads, feedback, enquiries, referrals"),
        ("Finance", "expenses, assets, tax, discounts"),
    ]),
]

for i, (icon, title, subtitle, badge_color, rows) in enumerate(arch_cols):
    x = Inches(0.4 + i * 4.25)
    y = Inches(1.35)
    w = Inches(4.05)
    h = Inches(5.6)
    shape = add_shape_bg(slide, x, y, w, h, CARD_BG, radius=0.06)

    # Header
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(0.45), Inches(0.4),
                 icon, font_size=20)
    add_text_box(slide, x + Inches(0.6), y + Inches(0.12), Inches(3.0), Inches(0.3),
                 title, font_size=17, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.6), y + Inches(0.4), Inches(3.0), Inches(0.25),
                 subtitle, font_size=10, color=GRAY)

    # Divider line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   x + Inches(0.2), y + Inches(0.75),
                                   w - Inches(0.4), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x2A, 0x3A, 0x4A)
    line.line.fill.background()

    # Rows
    row_y = y + Inches(0.9)
    r, g, b = badge_color[0], badge_color[1], badge_color[2]
    darker_badge = RGBColor(r // 4, g // 4, b // 4)
    lighter_text = RGBColor(min(r + 80, 255), min(g + 80, 255), min(b + 80, 255))
    for tech, desc in rows:
        # Badge
        bw = Inches(max(len(tech) * 0.09 + 0.25, 1.05))
        badge_shape = add_shape_bg(slide, x + Inches(0.2), row_y, bw, Inches(0.28),
                                    darker_badge, radius=0.2)
        add_text_box(slide, x + Inches(0.2), row_y + Inches(0.01), bw, Inches(0.26),
                     tech, font_size=9, color=lighter_text, bold=True, alignment=PP_ALIGN.CENTER)
        # Description
        add_text_box(slide, x + Inches(0.2) + bw + Inches(0.1), row_y + Inches(0.01),
                     w - bw - Inches(0.5), Inches(0.26),
                     desc, font_size=10, color=LIGHT_GRAY)
        row_y += Inches(0.62)

# Request flow bar at bottom
flow_y = Inches(7.05)
flow_items = [
    ("UI", "React + Ant Design"),
    ("API", "React Query"),
    ("Auth", "PyJWT + RBAC"),
    ("Route", "Flask Blueprint"),
    ("ORM", "MongoEngine"),
    ("DB", "MongoDB Atlas"),
]

# Background bar
add_shape_bg(slide, Inches(0.4), flow_y - Inches(0.05), SLIDE_W - Inches(0.8), Inches(0.48),
             RGBColor(0x0A, 0x10, 0x1E), radius=0.15)

fx = Inches(0.6)
for j, (label, tech) in enumerate(flow_items):
    add_text_box(slide, fx, flow_y - Inches(0.02), Inches(1.0), Inches(0.18),
                 label, font_size=7, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, fx, flow_y + Inches(0.14), Inches(1.6), Inches(0.2),
                 tech, font_size=10, color=WHITE, bold=True)
    fx += Inches(1.75)
    if j < len(flow_items) - 1:
        add_text_box(slide, fx - Inches(0.35), flow_y + Inches(0.05), Inches(0.3), Inches(0.3),
                     "\u2192", font_size=14, color=PRIMARY)

print("  Slide 9: Architecture")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 10: SECURITY
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_text_box(slide, Inches(0), Inches(0.5), SLIDE_W, Inches(0.6),
             "Enterprise-Grade Security", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.35),
             "Your business data is protected at every level", font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

sec_items = [
    ("\U0001F510", "JWT Authentication",
     "Secure token-based authentication with bcrypt\npassword hashing and session management\nfor every user in the system."),
    ("\U0001F465", "Role-Based Access Control",
     "Three-tier permissions: Owner (full system access),\nManager (branch-level), Staff (limited operations).\nEach role sees only what they need."),
    ("\U0001F3E2", "Branch Data Isolation",
     "Every record is branch-scoped. Staff see only\ntheir own branch data. Owners get full\ncross-branch visibility and control."),
    ("\U0001F4CA", "Complete Audit Trail",
     "Login history tracking, discount approval logs,\ncustomer merge audit trail, and deleted bill\narchives for full accountability."),
]

for i, (icon, title, desc) in enumerate(sec_items):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.9 + row * 2.6)
    w = Inches(5.6)
    h = Inches(2.2)
    shape = add_shape_bg(slide, x, y, w, h, CARD_BG, radius=0.06)

    # Icon circle
    icon_shape = add_shape_bg(slide, x + Inches(0.3), y + Inches(0.35),
                               Inches(0.65), Inches(0.65),
                               RGBColor(0x10, 0x3A, 0x36), radius=0.12)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.42), Inches(0.65), Inches(0.5),
                 icon, font_size=24, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, x + Inches(1.2), y + Inches(0.35), Inches(4.0), Inches(0.35),
                 title, font_size=18, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(1.2), y + Inches(0.8), Inches(4.0), Inches(1.1),
                 desc, font_size=13, color=GRAY, line_space=1.4)

print("  Slide 10: Security")

# ══════════════════════════════════════════════════════════════════════
# SLIDE 11: CTA / CLOSING
# ══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_text_box(slide, Inches(0), Inches(1.5), SLIDE_W, Inches(1.2),
             "Ready to Transform\nYour Salon Business?", font_size=48, color=WHITE,
             bold=True, alignment=PP_ALIGN.CENTER, line_space=1.15)

add_text_box(slide, Inches(2.5), Inches(3.0), Inches(8.333), Inches(0.7),
             "An all-in-one platform built for growth, efficiency,\nand exceptional customer experiences.",
             font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER, line_space=1.5)

# CTA button
btn_w, btn_h = Inches(3.5), Inches(0.7)
btn_x = (SLIDE_W - btn_w) // 2
btn = add_shape_bg(slide, btn_x, Inches(4.1), btn_w, btn_h, PRIMARY, radius=0.15)
add_text_box(slide, btn_x, Inches(4.15), btn_w, btn_h,
             "Get Started Today  \u2192", font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Feature pills
features = ["\u2713 Cloud-hosted", "\u2713 Mobile responsive", "\u2713 Multi-branch",
            "\u2713 GST compliant", "\u2713 WhatsApp integration"]
feat_x = Inches(1.5)
for feat in features:
    fw = Inches(len(feat) * 0.1 + 0.3)
    add_text_box(slide, feat_x, Inches(5.2), fw, Inches(0.35),
                 feat, font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)
    feat_x += fw + Inches(0.2)

add_text_box(slide, Inches(0), Inches(6.3), SLIDE_W, Inches(0.4),
             "Thank you for your time. Let\u2019s build something great together.",
             font_size=14, color=RGBColor(0x47, 0x55, 0x69), alignment=PP_ALIGN.CENTER)

print("  Slide 11: CTA")

# ══════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(__file__), "Saloon-Pro-Presentation.pptx")
prs.save(output_path)
print(f"\n  Saved: {output_path}")
print("  Done! Open in PowerPoint to present.")
